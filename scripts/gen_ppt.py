#!/usr/bin/env python3
"""Generate ScratchV promotion PPT."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
C_PRIMARY   = RGBColor(0x1A, 0x1A, 0x2E)  # dark navy
C_ACCENT    = RGBColor(0x00, 0x6D, 0x77)  # teal
C_ACCENT2   = RGBColor(0x2E, 0x86, 0xAB)  # lighter blue
C_ACCENT3   = RGBColor(0xF5, 0xA6, 0x23)  # warm orange
C_LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)  # light gray-blue
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TEXT  = RGBColor(0x2D, 0x2D, 0x2D)
C_GRAY      = RGBColor(0x66, 0x66, 0x66)
C_GREEN     = RGBColor(0x27, 0xAE, 0x60)

# Topic difficulty colors
DIFF_COLORS = {
    "低": RGBColor(0x27, 0xAE, 0x60),
    "中": RGBColor(0xF3, 0x9C, 0x12),
    "高": RGBColor(0xE7, 0x4C, 0x3C),
}

DIFF_TOPICS = {
    "低": [],
    "中": [],
    "高": [],
}

topics = [
    ("6",  "编译器性能测试套件",  "中"),
    ("7",  "编译器日志增强器",    "低"),
    ("9",  "DSL错误提示美化器",   "中"),
    ("1",  "DSL前端增强器",      "中"),
    ("13", "窥孔优化器",         "低"),
    ("14", "常量加载合并优化",    "低"),
    ("5",  "RISC-V汇编代码美化器","低"),
    ("20", "项目代码规范与格式化","低"),
    ("21", "IR 验证器",          "中"),
    ("28", "完善后端指令选择",     "中"),
    ("11", "控制流图（CFG）生成器","高"),
    ("12", "RISC-V后端指令计数统计器","高"),
    ("17", "寄存器分配（基本块内线性扫描）","高"),
    ("18", "指令调度（基本块内列表调度）","高"),
]

for n, name, diff in topics:
    DIFF_TOPICS[diff].append((n, name))

# ── Helper functions ──

def add_bg(slide, color=C_PRIMARY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, color=C_ACCENT):
    """Add a decorative left-bar shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(0.4), prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bottom_bar(slide, color=C_ACCENT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1),
        prs.slide_width, Inches(0.4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=C_DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.auto_size = None
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    tf.vertical_anchor = anchor
    return txBox


def add_rich_textbox(slide, left, top, width, height, lines,
                     font_name="Microsoft YaHei"):
    """lines: list of (text, size, color, bold, alignment, space_after)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.auto_size = None
    tf.word_wrap = True
    for i, item in enumerate(lines):
        text, size, color, bold, align, space = item
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(space)
    return txBox


def add_title_bar(slide, title_text, subtitle_text=""):
    add_bg(slide, C_WHITE)
    add_shape_bg(slide, C_ACCENT)
    # top accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0),
        Inches(12.929), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C_ACCENT
    line.line.fill.background()

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                title_text, font_size=32, color=C_PRIMARY, bold=True)
    if subtitle_text:
        add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                    subtitle_text, font_size=18, color=C_GRAY)


def add_bullet_card(slide, left, top, width, height, title, bullets,
                    title_color=C_ACCENT, bullet_color=C_DARK_TEXT):
    """Add a card with title and bullet points."""
    # card bg
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    card.line.width = Pt(1)
    card.shadow.inherit = False

    # title
    add_textbox(slide, left + Inches(0.3), top + Inches(0.15),
                width - Inches(0.6), Inches(0.5),
                title, font_size=18, color=title_color, bold=True)

    # bullets
    y = top + Inches(0.65)
    for b in bullets:
        txBox = slide.shapes.add_textbox(
            left + Inches(0.3), y, width - Inches(0.6), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"▸ {b}"
        p.font.size = Pt(14)
        p.font.color.rgb = bullet_color
        p.font.name = "Microsoft YaHei"
        y += Inches(0.35)


def add_page_number(slide, num, total):
    add_textbox(slide, Inches(12.0), Inches(7.1), Inches(1.2), Inches(0.4),
                f"{num}/{total}", font_size=10, color=C_WHITE,
                alignment=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, C_PRIMARY)

# Decorative shapes
for i in range(3):
    s = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(10 + i * 1.2), Inches(0.5 + i * 0.8),
        Inches(2.5), Inches(2.5)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3E)
    s.line.fill.background()
    s.fill.fore_color.brightness = 0.0

# Title
add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(1.2),
            "🧭 探索「AI模型 → 芯片指令」的\n神奇之旅",
            font_size=44, color=C_WHITE, bold=True)

# Subtitle
add_textbox(slide, Inches(1.2), Inches(3.0), Inches(10), Inches(0.8),
            "零基础友好的AI编译器开源项目 · 线上宣讲",
            font_size=24, color=C_ACCENT2)

# Bottom info
add_textbox(slide, Inches(1.2), Inches(5.0), Inches(6), Inches(0.5),
            "📅 2025年6月  |  🎯 三个月从零搭建你的第一个编译器",
            font_size=16, color=RGBColor(0xAA, 0xAA, 0xAA))

add_textbox(slide, Inches(1.2), Inches(5.5), Inches(6), Inches(0.5),
            "💡 不需要编译原理基础 · 不需要AI基础 · 只需要好奇心和耐心",
            font_size=14, color=RGBColor(0x88, 0x88, 0x88))


# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 — 目录
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "📋 宣讲提纲")

items = [
    ("01", "项目概览", "我们要做什么？"),
    ("02", "为什么值得参加", "你能收获什么"),
    ("03", "三个月学习路线", "Phase 1→2→3 逐级进阶"),
    ("04", "课题精选", "14个课题全方位解析"),
    ("05", "时间节点与适合人群", "关键里程碑 & 报名要求"),
    ("06", "Q&A", "常见疑问与报名"),
]
for i, (num, title, desc) in enumerate(items):
    row = i // 3
    col = i % 3
    x = Inches(0.8 + col * 4.2)
    y = Inches(2.0 + row * 2.5)
    # number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x, y + Inches(0.1), Inches(0.6), Inches(0.6)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = C_ACCENT
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = C_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Arial"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_textbox(slide, x + Inches(0.8), y, Inches(3.2), Inches(0.4),
                title, font_size=20, color=C_PRIMARY, bold=True)
    add_textbox(slide, x + Inches(0.8), y + Inches(0.4), Inches(3.2), Inches(0.4),
                desc, font_size=14, color=C_GRAY)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 — 你有没有好奇过
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "🤔 你有没有好奇过……", '一些可能让你“心里痒痒”的问题')

questions = [
    ("你写的 Python 代码，\n电脑到底怎么“听懂”并执行？",
     "从高级语言到机器指令，\n中间经历了怎样的魔法？"),
    ("那些炫酷的 AI 模型，\n最后怎么在小小的芯片上跑起来？",
     "模型是数学公式，\n芯片只懂 0 和 1，\n谁来当翻译官？"),
    ("编译器——\n这个听起来很高深的东西，\n到底在做什么？",
     "它不是一个黑盒，\n而是你可以亲手搭建的工具。"),
]

for i, (q, a) in enumerate(questions):
    x = Inches(0.8 + i * 4.2)
    y = Inches(2.0)
    # card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(3.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    card.line.width = Pt(1)

    add_textbox(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(1.5),
                q, font_size=18, color=C_PRIMARY, bold=True)
    add_textbox(slide, x + Inches(0.3), y + Inches(2.0), Inches(3.2), Inches(1.3),
                a, font_size=14, color=C_GRAY)

# punchline
add_textbox(slide, Inches(2), Inches(6.0), Inches(10), Inches(0.6),
            "—— 哪怕你只学过一点点编程，这个项目就是为你准备的 ——",
            font_size=20, color=C_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 — 项目概览
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "📌 我们要做什么？", "用三个月时间，从零搭建一个迷你AI编译器")

# Input → Process → Output
flow_data = [
    ("输入", "一个简单的 ONNX\nAI 模型文件\n（加法和乘法）", C_ACCENT),
    ("⬇\n编译器核心\n⬇", "读懂模型\n→ 中间语言\n→ 优化\n→ 生成指令", C_ACCENT2),
    ("输出", "RISC-V 汇编代码\n（add, load 等\n芯片指令）", C_ACCENT3),
]
for i, (title, desc, color) in enumerate(flow_data):
    x = Inches(0.8 + i * 4.2)
    y = Inches(2.2)
    # box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    add_textbox(slide, x + Inches(0.3), y + Inches(0.2), Inches(3.2), Inches(0.5),
                title, font_size=22, color=color, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.8), Inches(3.2), Inches(1.5),
                desc, font_size=16, color=C_DARK_TEXT,
                alignment=PP_ALIGN.CENTER)

# Key highlights
highlights = [
    "✅ 完全自主实现，不依赖 LLVM/MLIR 等巨型框架",
    "✅ 放到模拟器 (tinyfive) 中运行验证",
    "✅ 每一步都亲手写出来，真正搞懂原理",
]
for i, h in enumerate(highlights):
    add_textbox(slide, Inches(1.5), Inches(5.2 + i * 0.45), Inches(10), Inches(0.4),
                h, font_size=15, color=C_DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 — 技术全景图
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "🔬 从模型到芯片的完整链路", "编译器——AI 世界与芯片世界的桥梁")

# Pipeline
stages = [
    ("ONNX\n模型", "读取并解析\nAI 模型文件"),
    ("IR\n中间表示", "翻译成自己\n定义的指令"),
    ("优化器", "死代码消除\n常量折叠等"),
    ("后端", "指令选择\n寄存器分配"),
    ("RISC-V\n汇编", "生成芯片\n可执行指令"),
    ("模拟器\ntinyfive", "运行验证\n输出结果"),
]
for i, (title, desc) in enumerate(stages):
    x = Inches(0.5 + i * 2.1)
    y = Inches(2.5)
    # stage box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.9), Inches(2.8)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_PRIMARY if i % 2 == 0 else C_ACCENT
    shape.line.fill.background()

    add_textbox(slide, x + Inches(0.1), y + Inches(0.3), Inches(1.7), Inches(1.0),
                title, font_size=16, color=C_WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), y + Inches(1.5), Inches(1.7), Inches(1.0),
                desc, font_size=12, color=RGBColor(0xCC, 0xDD, 0xFF),
                alignment=PP_ALIGN.CENTER)

    if i < len(stages) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, x + Inches(1.9), y + Inches(1.2),
            Inches(0.2), Inches(0.3)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = C_ACCENT3
        arrow.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(5.8), Inches(12), Inches(0.5),
            "💡 整个流程没有任何黑盒——从模型文件的第一行，到汇编指令的最后一行，都由你亲手构建",
            font_size=16, color=C_GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 6 — 为什么值得参加
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "🔥 为什么值得你来试一试？")

reasons = [
    ("🧩 不需要大神基础", "学过一点 Python 或 C？够了。\n知道数组、函数、循环？完全够。\n两周带你入门 RISC-V，不需要预先掌握。",
     "像搭积木一样，一块一块搭起来", C_ACCENT),
    ("🛠️ 真正理解底层", "不再是只会 import torch 的调包侠。\n理解从数学模型到机器指令的完整链路。\n高性能计算、AI 芯片的核心能力。",
     "简历加分：独立实现AI→RISC-V完整编译器", C_ACCENT2),
    ("👥 温暖的开源社区", "每周线上答疑 + 讲解。\nPeer Review 代码，互相改 bug。\n一起庆祝每个里程碑的达成。",
     "你的代码会成为开源项目的一部分", C_ACCENT3),
]

for i, (title, desc, tag, color) in enumerate(reasons):
    x = Inches(0.5 + i * 4.2)
    y = Inches(2.0)

    # card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.9), Inches(4.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = color
    card.line.width = Pt(2)

    add_textbox(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.3), Inches(0.5),
                title, font_size=22, color=color, bold=True)
    add_textbox(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.3), Inches(2.0),
                desc, font_size=14, color=C_DARK_TEXT)
    # tag at bottom
    tag_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.3), y + Inches(3.8),
        Inches(3.3), Inches(0.5)
    )
    tag_shape.fill.solid()
    tag_shape.fill.fore_color.rgb = color
    tag_shape.line.fill.background()
    tf = tag_shape.text_frame
    tf.paragraphs[0].text = tag
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = C_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 — 三个月学习路线总览
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "🗺️ 三个月学习路线", "循序渐进的三阶段带飞计划")

phases = [
    ("Phase 1\n第 1~4 周", "跑通完整链路",
     [
         "跑通别人写好的完整例子",
         "看懂 模型→指令 的全过程",
         "建立对编译器的整体认知",
     ],
     "哇，原来是这样！", C_ACCENT),
    ("Phase 2\n第 5~9 周", "亲手实现编译器",
     [
         "把 ONNX 模型翻译成自己的 IR",
         "实现 RISC-V 后端生成汇编",
         "实现基础优化（常量折叠等）",
     ],
     "开始创造，成就感爆棚", C_ACCENT2),
    ("Phase 3\n第 10~12 周", "优化与产出",
     [
         "跑通更多模型",
         "指令更短、运行更快",
         "编写文档，完善项目",
     ],
     "我居然做出了一个编译器！", C_ACCENT3),
]

for i, (phase, subtitle, bullets, feeling, color) in enumerate(phases):
    x = Inches(0.5 + i * 4.2)
    y = Inches(2.0)

    # card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.9), Inches(4.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()

    # phase header
    add_textbox(slide, x + Inches(0.3), y + Inches(0.2), Inches(3.3), Inches(0.7),
                phase, font_size=18, color=C_WHITE, bold=True)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.9), Inches(3.3), Inches(0.4),
                subtitle, font_size=16, color=RGBColor(0xEE, 0xEE, 0xFF))

    # bullets (white card area)
    card2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.15), y + Inches(1.5),
        Inches(3.6), Inches(2.2)
    )
    card2.fill.solid()
    card2.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card2.line.fill.background()

    for j, b in enumerate(bullets):
        add_textbox(slide, x + Inches(0.4), y + Inches(1.7 + j * 0.6),
                    Inches(3.2), Inches(0.5),
                    f"✓ {b}", font_size=14, color=C_DARK_TEXT)

    # feeling
    add_textbox(slide, x + Inches(0.3), y + Inches(3.9), Inches(3.3), Inches(0.4),
                feeling, font_size=15, color=C_WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

# bottom note
add_textbox(slide, Inches(0.8), Inches(6.7), Inches(12), Inches(0.5),
            "⏱ 每周只需 8~10 小时 · 提供预置框架和 benchmark · 每周线上答疑",
            font_size=15, color=C_GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 8 — Phase 1 详解
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "📘 Phase 1 详解 —— 跑通完整链路 (第 1~4 周)",
              "目标：建立整体认知，看懂「模型→指令」的魔法")

details = [
    ("第 1 周", "环境搭建与入门",
     "安装工具链\n跑通第一个示例\n认识 ONNX 模型格式"),
    ("第 2 周", "RISC-V 基础",
     "学习 RISC-V 指令集基础\n理解寄存器、内存模型\n运行 tinyfive 模拟器"),
    ("第 3 周", "编译器工作流",
     "理解 前端→IR→后端 三段式\n跟着代码走一遍完整流程\n读懂中间表示的每一步转换"),
    ("第 4 周", "回顾与 checkpoint",
     "关键成果验收\n确保每个人都能讲清\n编译器的工作原理"),
]

for i, (week, title, desc) in enumerate(details):
    x = Inches(0.5 + i * 3.2)
    y = Inches(2.2)
    # week tag
    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.2), Inches(0.4)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = C_ACCENT
    tag.line.fill.background()
    tf = tag.text_frame
    tf.paragraphs[0].text = week
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = C_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, x, y + Inches(0.5), Inches(2.9), Inches(0.4),
                title, font_size=18, color=C_PRIMARY, bold=True)
    add_textbox(slide, x, y + Inches(1.0), Inches(2.9), Inches(2.0),
                desc, font_size=14, color=C_DARK_TEXT)

add_textbox(slide, Inches(0.8), Inches(6.5), Inches(12), Inches(0.4),
            "🎯 里程碑 checkpoint: 8 月 1 日 · 所有人能跑通完整链路",
            font_size=16, color=C_ACCENT3, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 — Phase 2 详解
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "📗 Phase 2 详解 —— 亲手实现编译器 (第 5~9 周)",
              "目标：自己动手实现核心功能")

phase2 = [
    ("DSL 前端", "▸ 设计并实现自己的 DSL\n▸ 解析 ONNX 模型\n▸ 输出自定义 IR\n\n课题：#1 DSL前端增强器\n课题：#9 错误提示美化器"),
    ("IR 与优化", "▸ 设计中间表示结构\n▸ 基础优化：常量折叠\n▸ 死代码消除\n\n课题：#21 IR 验证器\n课题：#13 窥孔优化器\n课题：#14 常量加载合并"),
    ("RISC-V 后端", "▸ 指令选择和映射\n▸ 寄存器分配初步\n▸ 生成可执行汇编\n\n课题：#28 后端指令选择\n课题：#17 寄存器分配"),
    ("集成验证", "▸ 打通前中后端\n▸ 在 tinyfive 上运行\n▸ 调试与排查问题\n\n课题：#12 指令计数统计\n课题：#6 性能测试套件"),
]

for i, (title, desc) in enumerate(phase2):
    x = Inches(0.5 + i * 3.2)
    y = Inches(2.2)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.9), Inches(4.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = C_ACCENT2
    card.line.width = Pt(1.5)

    add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.5), Inches(0.4),
                title, font_size=18, color=C_ACCENT2, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.7), Inches(2.5), Inches(3.3),
                desc, font_size=13, color=C_DARK_TEXT)

add_textbox(slide, Inches(0.8), Inches(6.6), Inches(12), Inches(0.4),
            "🎯 里程碑 checkpoint: 8 月 28 日 · 编译器跑通第一个完整模型",
            font_size=16, color=C_ACCENT3, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 — Phase 3 详解
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "📕 Phase 3 详解 —— 优化与产出 (第 10~12 周)",
              "目标：打磨作品，让编译器更聪明、更快")

phase3_items = [
    ("🚀 高级优化", [
        "指令调度——让汇编指令排列更高效",
        "寄存器分配——线性扫描算法",
        "控制流图（CFG）分析与优化",
    ]),
    ("📊 评测与完善", [
        "编译器性能测试与基准评测",
        "RISC-V 指令计数统计分析",
        "跨多个模型的正确性验证",
    ]),
    ("📝 文档与开源", [
        "编写项目文档和使用指南",
        "代码规范与格式化统一",
        "开源社区贡献指南",
    ]),
    ("🎓 结项展示", [
        "成果演示 + 技术分享",
        "每人总结自己的贡献与收获",
        "项目成为简历上的亮点",
    ]),
]

for i, (title, items) in enumerate(phase3_items):
    x = Inches(0.5 + i * 3.2)
    y = Inches(2.2)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.9), Inches(3.8)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = C_PRIMARY if i % 2 == 0 else C_ACCENT
    card.line.fill.background()

    add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.5), Inches(0.5),
                title, font_size=18, color=C_WHITE, bold=True)
    for j, item in enumerate(items):
        add_textbox(slide, x + Inches(0.2), y + Inches(0.8 + j * 0.6),
                    Inches(2.5), Inches(0.5),
                    f"▸ {item}", font_size=14,
                    color=RGBColor(0xE8, 0xEE, 0xFF))

add_textbox(slide, Inches(0.8), Inches(6.3), Inches(12), Inches(0.6),
            "🎓 结项: 9 月 27 日 · 成果展示 + 结项总结 · 项目成为你简历上的骄傲",
            font_size=16, color=C_ACCENT3, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(6.8), Inches(12), Inches(0.4),
            "课题相关: #11 CFG生成器 · #12 指令计数统计 · #17 寄存器分配 · #18 指令调度 · #5 汇编美化器 · #6 性能测试套件",
            font_size=12, color=C_GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 11 — 课题精选总览
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "🎯 课题精选总览", "14 个课题覆盖编译器的各个模块，按难度分层")

# Table header
header_y = Inches(2.2)
col_x = [Inches(0.8), Inches(1.8), Inches(5.5)]
col_w = [Inches(0.8), Inches(3.5), Inches(1.5)]

for ci, (txt, w) in enumerate([("编号", 0.8), ("课题名称", 10), ("难度", 1.5)]):
    add_textbox(slide, col_x[ci], header_y, Inches(w), Inches(0.4),
                txt, font_size=13, color=C_ACCENT, bold=True)

# divider
div = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.6), Inches(11.5), Inches(0.02)
)
div.fill.solid()
div.fill.fore_color.rgb = C_ACCENT
div.line.fill.background()

for i, (num, name, diff) in enumerate(topics):
    y = Inches(2.7 + i * 0.32)
    add_textbox(slide, col_x[0], y, Inches(0.8), Inches(0.3),
                num, font_size=12, color=C_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.0), y, Inches(8.5), Inches(0.3),
                name, font_size=13, color=C_DARK_TEXT)

    # diff badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.2), y - Inches(0.02),
        Inches(0.8), Inches(0.3)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = DIFF_COLORS[diff]
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].text = diff
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = C_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER


# ══════════════════════════════════════════════════════════════════════
# SLIDE 12 — 课题详解：低难度
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "🌟 课题详解 · 入门友好型 (难度：低)", "适合新手起步，快速获得成就感")

entries_low = [
    ("#7 编译器日志增强器", "为编译器添加结构化日志输出\n支持不同日志级别（INFO/WARN/ERROR）\n帮助开发者调试编译流程"),
    ("#13 窥孔优化器", "扫描生成的汇编代码\n用更短的指令序列替换冗余模式\n经典的局部优化技术"),
    ("#14 常量加载合并", "识别重复的常量加载指令\n合并为一次加载，减少代码体积\n简单但效果明显的优化"),
    ("#5 RISC-V 汇编美化器", "格式化汇编输出\n添加注释和可读性改进\n让生成的代码更易读"),
    ("#20 项目代码规范与格式化", "统一代码风格\n配置 linter/formatter\n提升项目的专业性"),
]

for i, (title, desc) in enumerate(entries_low):
    x = Inches(0.5 + (i % 3) * 4.2)
    y = Inches(2.0 + (i // 3) * 2.6)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.9), Inches(2.3)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = C_GREEN
    card.line.width = Pt(1.5)

    # diff badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.2), y + Inches(0.15),
        Inches(0.6), Inches(0.3)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = C_GREEN
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].text = "低"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = C_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, x + Inches(1.0), y + Inches(0.15), Inches(2.7), Inches(0.4),
                title, font_size=16, color=C_GREEN, bold=True)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.6), Inches(3.4), Inches(1.5),
                desc, font_size=13, color=C_DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 13 — 课题详解：中等难度
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "⚡ 课题详解 · 进阶挑战型 (难度：中)", "需要一定的编程能力，深入编译器核心")

mid_topics = [
    ("#1 DSL 前端增强器",
     "扩展 DSL 语法支持\n支持更多算子类型\n改进前端解析能力"),
    ("#9 DSL 错误提示美化器",
     "友好的错误信息展示\n精确的错误位置定位\n彩色终端输出"),
    ("#6 编译器性能测试套件",
     "自动化 benchmark 流程\n统计编译时间和代码质量\n生成性能对比报告"),
]

for i, (title, desc) in enumerate(mid_topics):
    x = Inches(0.5 + i * 4.2)
    y = Inches(2.2)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.9), Inches(2.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = DIFF_COLORS["中"]
    card.line.width = Pt(1.5)

    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.2), y + Inches(0.15),
        Inches(0.6), Inches(0.3)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = DIFF_COLORS["中"]
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].text = "中"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = C_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, x + Inches(1.0), y + Inches(0.15), Inches(2.7), Inches(0.4),
                title, font_size=16, color=DIFF_COLORS["中"], bold=True)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.65), Inches(3.4), Inches(1.5),
                desc, font_size=13, color=C_DARK_TEXT)

# More mid topics in second row
mid2 = [
    ("#21 IR 验证器", "检查IR的合法性\n验证数据流正确性\n确保优化不改变语义"),
    ("#28 完善后端指令选择", "扩展指令选择规则\n支持更多 RISC-V 指令\n后端功能完善"),
]
for i, (title, desc) in enumerate(mid2):
    x = Inches(2.6 + i * 4.2)
    y = Inches(5.0)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.9), Inches(2.0)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = DIFF_COLORS["中"]
    card.line.width = Pt(1.5)

    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.2), y + Inches(0.15),
        Inches(0.6), Inches(0.3)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = DIFF_COLORS["中"]
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].text = "中"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = C_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, x + Inches(1.0), y + Inches(0.15), Inches(2.7), Inches(0.4),
                title, font_size=16, color=DIFF_COLORS["中"], bold=True)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.65), Inches(3.4), Inches(1.0),
                desc, font_size=13, color=C_DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 14 — 课题详解：高难度
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "🏆 课题详解 · 大神修炼型 (难度：高)", "挑战编译器核心算法，收获满满硬技能")

high_topics = [
    ("#11 控制流图（CFG）生成器",
     "从 IR 中构建基本块\n生成控制流图可视化\n为后续优化奠定基础",
     "编译优化 · 图算法"),
    ("#12 RISC-V 后端指令计数统计器",
     "统计生成的各类指令数量\n分析指令分布特征\n输出可视化统计报告",
     "数据分析 · 可视化"),
    ("#17 寄存器分配（基本块内线性扫描）",
     "实现寄存器分配算法\n处理变量生命周期\n优化寄存器使用效率",
     "经典算法实现"),
    ("#18 指令调度（基本块内列表调度）",
     "分析指令依赖关系\n重排指令顺序\n提高流水线效率",
     "编译优化 · 调度算法"),
]

for i, (title, desc, tag) in enumerate(high_topics):
    x = Inches(0.5 + i * 3.2)
    y = Inches(2.2)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.9), Inches(4.0)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xF5, 0xF5)
    card.line.color.rgb = DIFF_COLORS["高"]
    card.line.width = Pt(2)

    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.2), y + Inches(0.15),
        Inches(0.6), Inches(0.3)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = DIFF_COLORS["高"]
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].text = "高"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = C_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, x + Inches(0.2), y + Inches(0.6), Inches(2.5), Inches(0.5),
                title, font_size=14, color=DIFF_COLORS["高"], bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.2), Inches(2.5), Inches(1.8),
                desc, font_size=12, color=C_DARK_TEXT)

    # tag
    tag_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.2), y + Inches(3.5),
        Inches(2.5), Inches(0.35)
    )
    tag_shape.fill.solid()
    tag_shape.fill.fore_color.rgb = DIFF_COLORS["高"]
    tag_shape.line.fill.background()
    tf = tag_shape.text_frame
    tf.paragraphs[0].text = tag
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = C_WHITE
    tf.paragraphs[0].font.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

add_textbox(slide, Inches(0.8), Inches(6.5), Inches(12), Inches(0.5),
            "💡 高难度课题有 mentor 全程指导，不用担心做不出来",
            font_size=15, color=C_GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 15 — 适合人群
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "🙋 谁适合报名？")

checks = [
    "大二、大三、研一，或者自学编程爱好者",
    "学过一门编程语言（Python / C / C++ 都行）",
    "对“计算机到底怎么跑程序”有好奇心",
    "不怕犯错，敢写代码（Bug 是学习的一部分！）",
    "每周能拿出 8～10 小时",
]

for i, c in enumerate(checks):
    row = i // 3
    col = i % 3
    x = Inches(0.8 + col * 4.2)
    y = Inches(2.2 + row * 1.2)
    check_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(0.8)
    )
    check_shape.fill.solid()
    check_shape.fill.fore_color.rgb = RGBColor(0xF0, 0xFF, 0xF0)
    check_shape.line.color.rgb = C_GREEN
    check_shape.line.width = Pt(1)

    add_textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(0.5),
                f"✅  {c}", font_size=15, color=C_DARK_TEXT)

# you don't need
add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.4),
            "你可能还没学过编译原理、还没搞懂指令集、甚至对汇编有点畏惧——都没关系。",
            font_size=16, color=C_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.4),
            "我们就是来带你一步步跨过这些坎的。",
            font_size=18, color=C_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 16 — 时间节点
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "📅 关键时间节点")

timeline = [
    ("即日起", "开始报名"),
    ("6 月 20 日", "线上宣讲 + 课题选择\n报名截止"),
    ("7 月 10 日", "项目正式开启\n启动会 + 第一周任务"),
    ("8 月 1 日", "Phase 1 关键成果验收\n第一个里程碑"),
    ("8 月 28 日", "Phase 2 关键成果验收\n第二个里程碑"),
    ("9 月 27 日", "项目结项\n成果展示 + 总结"),
]

for i, (date, event) in enumerate(timeline):
    x = Inches(0.8 + i * 2.1)
    y = Inches(2.5)

    # dot
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(0.5), y, Inches(0.3), Inches(0.3)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = C_ACCENT if i < len(timeline) - 1 else C_ACCENT3
    dot.line.fill.background()

    # connector line
    if i < len(timeline) - 1:
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x + Inches(0.8), y + Inches(0.12),
            Inches(1.3), Inches(0.06)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = C_ACCENT
        line.line.fill.background()

    # date
    add_textbox(slide, x, y + Inches(0.4), Inches(1.9), Inches(0.4),
                date, font_size=14, color=C_PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)
    # event
    add_textbox(slide, x, y + Inches(0.8), Inches(1.9), Inches(1.2),
                event, font_size=12, color=C_DARK_TEXT,
                alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 17 — FAQ
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "💬 常见疑问")

faqs = [
    ("没听过 ONNX，能行吗？",
     "当然可以。第 1 周就会带你跑通例子，ONNX 只是一个文件格式，把它当成“模型存盘”就好。"),
    ("没学过编译原理，会不会听不懂？",
     "我们会用直观的比喻（编译器就像“翻译官”），避开理论轰炸，先动手再做总结。"),
    ("需要买 RISC-V 开发板吗？",
     "不需要。全程用软件模拟器（tinyfive），在笔记本电脑上就能跑。"),
    ("中途跟不上怎么办？",
     "每个阶段有进度检查，mentor 会主动帮忙。有缓冲时间，可选只完成核心路径。完成比完美重要。"),
]

for i, (q, a) in enumerate(faqs):
    y = Inches(2.0 + i * 1.3)
    q_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.5), Inches(0.45)
    )
    q_bg.fill.solid()
    q_bg.fill.fore_color.rgb = C_ACCENT
    q_bg.line.fill.background()

    add_textbox(slide, Inches(1.0), y + Inches(0.03), Inches(11), Inches(0.4),
                f"Q: {q}", font_size=15, color=C_WHITE, bold=True)

    add_textbox(slide, Inches(1.0), y + Inches(0.55), Inches(11), Inches(0.6),
                f"A: {a}", font_size=14, color=C_DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 18 — 报名信息
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_PRIMARY)

# decorative
for i in range(4):
    s = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(10 + i * 1.5), Inches(5 - i * 0.8),
        Inches(2), Inches(2)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3E)
    s.line.fill.background()

add_textbox(slide, Inches(1.2), Inches(1.0), Inches(11), Inches(1.2),
            "🌟 从今天起，给自己一个\n“创造编译器”的机会",
            font_size=40, color=C_WHITE, bold=True)

add_textbox(slide, Inches(1.2), Inches(3.0), Inches(11), Inches(0.6),
            "也许你现在觉得编译器遥不可及，\n但三个月后，你会看着自己写的代码，把一行行模型规则变成芯片指令。",
            font_size=18, color=RGBColor(0xCC, 0xCC, 0xDD))

add_textbox(slide, Inches(1.2), Inches(4.0), Inches(11), Inches(0.5),
            "“我居然做到了”——这会是大学期间最难忘的回忆之一。",
            font_size=18, color=C_ACCENT3, bold=True)

add_textbox(slide, Inches(1.2), Inches(5.0), Inches(11), Inches(0.4),
            "不要让“基础不够”成为不敢开始的理由。",
            font_size=20, color=C_WHITE, bold=True)

add_textbox(slide, Inches(1.2), Inches(5.7), Inches(11), Inches(0.4),
            "👉 立即报名：https://your-form-link.com      📧 咨询：mentor@example.com",
            font_size=16, color=RGBColor(0xAA, 0xCC, 0xFF))

add_textbox(slide, Inches(1.2), Inches(6.3), Inches(11), Inches(0.4),
            "QQ 群：xxxxxxxxx    欢迎转发给同样好奇的小伙伴！",
            font_size=14, color=RGBColor(0x88, 0x88, 0xAA))

add_textbox(slide, Inches(1.2), Inches(6.8), Inches(11), Inches(0.3),
            "你不需要很厉害才能开始，但你需要开始才能很厉害。",
            font_size=14, color=C_ACCENT3, bold=False, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 19 — Thank you
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_PRIMARY)

add_textbox(slide, Inches(1.2), Inches(2.0), Inches(11), Inches(1.5),
            "Thank You 🙌",
            font_size=56, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.2), Inches(3.8), Inches(11), Inches(0.8),
            "期待与你一起，开启编译器的奇妙旅程",
            font_size=24, color=C_ACCENT2, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.2), Inches(5.0), Inches(11), Inches(0.5),
            "Questions & Answers",
            font_size=20, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.CENTER)

# ── Save ──
output_path = "/home/kinsomwang/workspace/ScratchV/ScratchV_宣讲PPT.pptx"
prs.save(output_path)
print(f"✅ PPT saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
