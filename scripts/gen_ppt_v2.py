"""
ScratchV — Professional Project Promotion PPT
Design: modern, clean, high-end open-source style
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

prs = Presentation()
W = prs.slide_width = Inches(13.333)
H = prs.slide_height = Inches(7.5)

# ── Color System ──
C_BG_DARK    = RGBColor(0x0A, 0x0E, 0x1A)  # deep dark
C_BG_CARD    = RGBColor(0x12, 0x18, 0x2A)  # card dark
C_BG_LIGHT   = RGBColor(0xF8, 0xFA, 0xFC)  # light bg
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT_PRIMARY  = RGBColor(0xE8, 0xEC, 0xF4)  # light text
C_TEXT_SECONDARY = RGBColor(0x94, 0xA3, 0xB8)  # muted text
C_TEXT_DARK   = RGBColor(0x1E, 0x29, 0x3B)  # dark text for light cards

# Brand colors
C_TEAL       = RGBColor(0x2D, 0xD4, 0xBF)  # primary brand - teal
C_TEAL_DIM   = RGBColor(0x14, 0x8A, 0x7A)  # darker teal
C_PURPLE     = RGBColor(0x8B, 0x5C, 0xF6)  # accent purple
C_PURPLE_DIM = RGBColor(0x6D, 0x3A, 0xD6)  # darker purple
C_ORANGE     = RGBColor(0xF5, 0x9E, 0x0B)  # warm accent
C_RED        = RGBColor(0xEF, 0x44, 0x44)
C_GREEN      = RGBColor(0x10, 0xB9, 0x81)

# Gradients (simulated with shapes)
ACCENT_GRADIENT = [C_TEAL, C_PURPLE]

# ── Topic data ──
TOPICS = [
    ("#1",  "DSL Frontend Enhancer",          "中", "Extend DSL syntax, support more operators"),
    ("#5",  "RISC-V Assembly Beautifier",      "低", "Format assembly output, add readability"),
    ("#6",  "Compiler Performance Bench",      "中", "Automated benchmark & report"),
    ("#7",  "Compiler Log Enhancer",           "低", "Structured logging with levels"),
    ("#9",  "DSL Error Message Beautifier",    "中", "Friendly error display & location"),
    ("#11", "CFG Generator",                   "高", "Build basic blocks, visualize CFG"),
    ("#12", "RISC-V Instr. Counter",           "高", "Count & analyze instruction distribution"),
    ("#13", "Peephole Optimizer",              "低", "Replace redundant instruction patterns"),
    ("#14", "Const Load Merging",              "低", "Merge duplicate constant loads"),
    ("#17", "Register Allocator (L.Scan)",     "高", "Linear scan register allocation"),
    ("#18", "Instr. Scheduler (L.List)",       "高", "Dependency analysis & rescheduling"),
    ("#20", "Code Style & Formatting",         "低", "Unified code style & linter config"),
    ("#21", "IR Verifier",                     "中", "Validate IR legality & data flow"),
    ("#28", "Backend Instr. Selection",        "中", "Expand RISC-V instruction mapping"),
]

DIFF_LABELS = {"低": "Beginner", "中": "Intermediate", "高": "Advanced"}
DIFF_COLORS = {"低": C_GREEN, "中": C_ORANGE, "高": C_RED}

# ── Helper functions ──

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_bar(slide, left, top, width, height, color1, color2):
    """Simulate gradient with overlapping shapes."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color1
    shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color, radius=None,
             border_color=None, border_width=None, opacity=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    if radius:
        # Adjust corner radius via XML
        spPr = shape._element.find(qn('a:spPr')) if shape._element.find(qn('a:spPr')) is not None else None
        if spPr is None:
            spPr = shape._element.find(qn('p:spPr'))
        if spPr is not None and radius:
            # Round each corner
            for attr in ['b', 'l', 't', 'r']:
                pass  # python-pptx doesn't easily support custom corner radius
    if opacity is not None and 0 <= opacity <= 1:
        shape.fill.fore_color.brightness = opacity
    return shape


def add_text(slide, left, top, width, height, text, size=16,
             color=C_TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT,
             font="Inter", anchor=MSO_ANCHOR.TOP, spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.auto_size = None
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.space_after = Pt(spacing or 0)
    return txBox


def add_rich_text(slide, left, top, width, height, lines, font="Inter"):
    """lines: [(text, size, color, bold, align, space_after), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.auto_size = None
    tf.word_wrap = True
    for i, (text, size, color, bold, align, space) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
        p.space_after = Pt(space)
    return txBox


def add_tag(slide, left, top, width, height, text, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = C_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Inter"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def slide_number(slide, num, total):
    add_text(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.35),
             f"{num:02d}  /  {total:02d}", size=9,
             color=C_TEXT_SECONDARY, align=PP_ALIGN.RIGHT)


def add_section_header(slide, section_num, title, subtitle,
                       is_dark=True, num_total=None):
    bg = C_BG_DARK if is_dark else C_BG_LIGHT
    set_slide_bg(slide, bg)

    # Top accent line
    add_gradient_bar(slide, Inches(0), Inches(0), W, Inches(0.04),
                     C_TEAL, C_PURPLE)

    # Section number
    add_text(slide, Inches(0.8), Inches(0.6), Inches(1), Inches(0.5),
             f"/ {section_num:02d}", size=14, color=C_TEAL, bold=True,
             font="JetBrains Mono")

    # Title
    c_title = C_TEXT_PRIMARY if is_dark else C_TEXT_DARK
    add_text(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(0.8),
             title, size=36, color=c_title, bold=True, spacing=4)

    # Subtitle
    c_sub = C_TEXT_SECONDARY if is_dark else C_GRAY
    add_text(slide, Inches(0.8), Inches(1.8), Inches(10), Inches(0.5),
             subtitle, size=16, color=c_sub, spacing=2)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)

# Decorative gradient shapes
for i, (x, y, s, c) in enumerate([
    (10.5, -1.5, 5, C_TEAL), (8.5, 5.0, 4, C_PURPLE), (-1, 4, 3.5, C_TEAL_DIM)
]):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                                    Inches(s), Inches(s))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c
    shape.fill.fore_color.brightness = 0.85
    shape.line.fill.background()

# Top accent bar
add_gradient_bar(slide, Inches(0), Inches(0), W, Inches(0.06), C_TEAL, C_PURPLE)

# Logo / brand
add_text(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
         "SCRATCHV", size=12, color=C_TEAL, bold=True,
         font="JetBrains Mono", spacing=4)

# Main title
add_text(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(1.5),
         "From AI Models to Silicon Instructions:\nBuilding Your First Compiler",
         size=48, color=C_WHITE, bold=True, spacing=6)

# Subtitle
add_text(slide, Inches(0.8), Inches(3.8), Inches(8), Inches(0.6),
         "An open-source, beginner-friendly journey into the world of compilers, "
         "AI acceleration, and RISC-V.",
         size=18, color=C_TEXT_SECONDARY, spacing=4)

# Tagline box
box = add_rect(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(0.55),
               C_BG_CARD, radius=True, border_color=C_TEAL, border_width=Pt(0.5))
add_text(slide, Inches(1.0), Inches(5.05), Inches(5.2), Inches(0.45),
         "  \U0001f680  Open Source  ·  Hands-on  ·  12 Weeks  ·  Mentor-Guided",
         size=13, color=C_TEXT_SECONDARY, align=PP_ALIGN.LEFT)

# Right side info
add_text(slide, Inches(9.5), Inches(5.5), Inches(3.5), Inches(0.4),
         "Project Launch: July 2025", size=11,
         color=C_TEXT_SECONDARY, align=PP_ALIGN.RIGHT)
add_text(slide, Inches(9.5), Inches(5.85), Inches(3.5), Inches(0.4),
         "Target: RISC-V (via tinyfive sim)", size=11,
         color=C_TEXT_SECONDARY, align=PP_ALIGN.RIGHT)

add_text(slide, Inches(0.8), Inches(6.8), Inches(5), Inches(0.4),
         "#ScratchV  #Compiler  #RISC-V  #OpenSource  #AI",
         size=10, color=C_TEXT_SECONDARY)

slide_number(slide, 1, 19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Vision
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 1, "The Vision",
                   "Three questions that sparked this project", num_total=19)

cards = [
    ("01", "How does Python\nbecome machine code?",
     "From your editor to the CPU — what happens\nin between? Compilers are the hidden\nbridge we use every day.",
     C_TEAL),
    ("02", "How do AI models\nrun on tiny chips?",
     "Neural networks are math; chips speak\n0s and 1s. Who translates? That’s the\ncompiler’s superpower.",
     C_PURPLE),
    ("03", "Can I build\na compiler myself?",
     "A compiler isn’t magic — it’s just a program.\nIn 12 weeks, you’ll write one from scratch\nand understand every single line.",
     C_ORANGE),
]

for i, (num, title, desc, color) in enumerate(cards):
    x = Inches(0.6 + i * 4.2)
    y = Inches(2.6)

    # Number
    add_text(slide, x, y, Inches(0.6), Inches(0.5), num,
             size=24, color=color, bold=True, font="JetBrains Mono", spacing=2)

    # Card bg
    card = add_rect(slide, x, y + Inches(0.5), Inches(3.8), Inches(3.2),
                    C_BG_CARD, radius=True, border_color=color, border_width=Pt(1))

    # Title
    add_text(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.2), Inches(1.0),
             title, size=20, color=C_WHITE, bold=True, spacing=4)

    # Desc
    add_text(slide, x + Inches(0.3), y + Inches(1.9), Inches(3.2), Inches(1.5),
             desc, size=13, color=C_TEXT_SECONDARY, spacing=4)

# Bottom note
add_text(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
         "You don’t need a CS degree.  You don’t need to know RISC-V.  "
         "You just need curiosity and 8–10 hours a week.",
         size=15, color=C_TEAL, bold=False, align=PP_ALIGN.CENTER, spacing=2)

slide_number(slide, 2, 19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 — What We Build
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 2, "What We Build",
                   "A minimal AI compiler: ONNX model → RISC-V assembly, entirely from scratch")

# Pipeline
stages = [
    ("ONNX\nModel", "Read & parse\nthe AI model\nfile format", C_TEAL),
    ("Custom\nIR", "Translate into\nyour very own\nintermediate rep.", C_PURPLE),
    ("Optimizer", "Dead code elim.\nconstant folding\npeephole opt.", C_TEAL_DIM),
    ("RISC-V\nBackend", "Instruction selection\nregister allocation\ncode emission", C_PURPLE_DIM),
    ("tinyfive\nSimulator", "Run & verify\non a simulated\nRISC-V chip", C_ORANGE),
]

for i, (title, desc, color) in enumerate(stages):
    x = Inches(0.4 + i * 2.55)
    y = Inches(2.8)

    # Stage card
    card = add_rect(slide, x, y, Inches(2.3), Inches(3.0),
                    C_BG_CARD, radius=True, border_color=color, border_width=Pt(1))

    add_text(slide, x + Inches(0.2), y + Inches(0.3), Inches(1.9), Inches(0.8),
             title, size=16, color=color, bold=True, align=PP_ALIGN.CENTER,
             spacing=4)
    add_text(slide, x + Inches(0.2), y + Inches(1.3), Inches(1.9), Inches(1.5),
             desc, size=12, color=C_TEXT_SECONDARY, align=PP_ALIGN.CENTER,
             spacing=3)

    # Arrow between stages
    if i < len(stages) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            x + Inches(2.3), y + Inches(1.3),
            Inches(0.25), Inches(0.25)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = color
        arrow.line.fill.background()

# Highlights
highlights_data = [
    ("Zero black boxes", "Every line is written by you"),
    ("No LLVM/MLIR", "Understand without framework crutches"),
    ("Real execution", "Run on tinyfive simulator"),
]

for i, (h, d) in enumerate(highlights_data):
    x = Inches(0.6 + i * 4.2)
    y = Inches(6.1)
    add_text(slide, x, y, Inches(3.8), Inches(0.3),
             f"✓  {h}", size=13, color=C_GREEN, bold=True, spacing=2)
    add_text(slide, x, y + Inches(0.3), Inches(3.8), Inches(0.3),
             d, size=11, color=C_TEXT_SECONDARY, spacing=2)

slide_number(slide, 3, 19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 — Why Participate
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 3, "Why Join?",
                   "Three reasons this project is different")

reasons = [
    ("\U0001f9f1", "Build Real Understanding",
     [
         "Not another “import torch” tutorial",
         "From math formula to machine instruction",
         "Core skill for HPC, AI chips, systems",
         "Resume line: built a full AI→RISC-V compiler",
     ]),
    ("\U0001f30d", "Warm Community",
     [
         "Weekly live Q&A and deep-dive sessions",
         "Peer code review — learn by teaching",
         "Mentor-guided, no one left behind",
         "Your code helps future learners",
     ]),
    ("\U0001f680", "Low Floor, High Ceiling",
     [
         "Only need basic Python or C",
         "14 topics across 3 difficulty tiers",
         "Start with a beautifier, end with a register allocator",
         "“Done is better than perfect” — flexible pacing",
     ]),
]

for i, (emoji, title, bullets) in enumerate(reasons):
    x = Inches(0.5 + i * 4.2)
    y = Inches(2.6)

    card = add_rect(slide, x, y, Inches(3.9), Inches(4.0),
                    C_BG_CARD, radius=True, border_color=C_TEAL, border_width=Pt(0.5))

    add_text(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.3), Inches(0.5),
             f"{emoji}  {title}", size=20, color=C_WHITE, bold=True, spacing=4)

    for j, b in enumerate(bullets):
        add_text(slide, x + Inches(0.3), y + Inches(1.1 + j * 0.55),
                 Inches(3.3), Inches(0.5),
                 f"▸  {b}", size=13, color=C_TEXT_SECONDARY, spacing=3)

slide_number(slide, 4, 19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 — Timeline Overview
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 4, "12-Week Roadmap",
                   "Three phases. Clear milestones. Measurable progress.")

phases = [
    ("Phase 1", "Weeks 1–4", "Foundation",
     [
         "Set up toolchain & run first example",
         "Learn RISC-V basics & tinyfive sim",
         "Walk through compiler end-to-end",
         "Milestone: understand the full pipeline",
     ], C_TEAL, "\U0001f4d6"),
    ("Phase 2", "Weeks 5–9", "Implementation",
     [
         "Parse ONNX into custom IR",
         "Build RISC-V code generator backend",
         "Implement basic optimizations",
         "Milestone: first model runs on tinyfive",
     ], C_PURPLE, "\U0001f528"),
    ("Phase 3", "Weeks 10–12", "Optimization",
     [
         "Advanced opt: CFG, reg alloc, scheduling",
         "Performance benchmarking & tuning",
         "Documentation & open-source polish",
         "Milestone: final demo & project wrap-up",
     ], C_ORANGE, "\U0001f3c6"),
]

for i, (phase, period, title, bullets, color, emoji) in enumerate(phases):
    x = Inches(0.5 + i * 4.2)
    y = Inches(2.6)

    # Phase header
    header = add_rect(slide, x, y, Inches(3.9), Inches(1.0),
                      color, radius=False)
    add_text(slide, x + Inches(0.3), y + Inches(0.1), Inches(2.5), Inches(0.4),
             f"{emoji}  {phase}", size=16, color=C_WHITE, bold=True, spacing=2)
    add_text(slide, x + Inches(0.3), y + Inches(0.5), Inches(2.5), Inches(0.4),
             period, size=11, color=C_WHITE, spacing=2)

    # Card body
    card = add_rect(slide, x, y + Inches(1.0), Inches(3.9), Inches(3.4),
                    C_BG_CARD, radius=False)
    # Round bottom corners
    add_text(slide, x + Inches(0.3), y + Inches(1.1), Inches(3.3), Inches(0.4),
             title, size=18, color=color, bold=True, spacing=4)

    for j, b in enumerate(bullets):
        add_text(slide, x + Inches(0.3), y + Inches(1.7 + j * 0.5),
                 Inches(3.3), Inches(0.5),
                 f"✓  {b}", size=12, color=C_TEXT_SECONDARY, spacing=3)

    # Time commitment
    add_text(slide, x + Inches(0.3), y + Inches(3.8), Inches(3.3), Inches(0.3),
             "⏱  8–10 hrs/week", size=10, color=C_TEXT_SECONDARY, spacing=2)

slide_number(slide, 5, 19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 6 — Curriculum Detail
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 5, "Curriculum at a Glance",
                   "What you’ll learn, week by week")

# Week grid
weeks = [
    ("W1", "Toolchain Setup", C_TEAL),
    ("W2", "RISC-V Basics", C_TEAL),
    ("W3", "Compiler Pipeline", C_TEAL),
    ("W4", "Phase-1 Checkpoint", C_TEAL_DIM),
    ("W5", "DSL Frontend", C_PURPLE),
    ("W6", "IR Design", C_PURPLE),
    ("W7", "Backend I", C_PURPLE),
    ("W8", "Optimization", C_PURPLE),
    ("W9", "Phase-2 Checkpoint", C_PURPLE_DIM),
    ("W10", "Adv. Optimization", C_ORANGE),
    ("W11", "Benchmark & Polish", C_ORANGE),
    ("W12", "Final Demo", C_RED),
]

for i, (week, topic, color) in enumerate(weeks):
    x = Inches(0.5 + (i % 6) * 2.1)
    y = Inches(2.6 + (i // 6) * 1.4)

    card = add_rect(slide, x, y, Inches(1.9), Inches(1.1),
                    C_BG_CARD, radius=True, border_color=color, border_width=Pt(0.5))

    add_text(slide, x + Inches(0.1), y + Inches(0.1), Inches(1.7), Inches(0.3),
             week, size=10, color=color, bold=True, align=PP_ALIGN.CENTER,
             font="JetBrains Mono", spacing=2)
    add_text(slide, x + Inches(0.1), y + Inches(0.4), Inches(1.7), Inches(0.5),
             topic, size=12, color=C_TEXT_PRIMARY, align=PP_ALIGN.CENTER, spacing=2)

# Milestones
add_text(slide, Inches(0.5), Inches(5.8), Inches(12), Inches(0.4),
         "\U0001f4cc  Milestones:  Aug 1 (Phase 1)  ·  Aug 28 (Phase 2)  ·  Sep 27 (Final Demo)",
         size=14, color=C_TEAL, align=PP_ALIGN.CENTER, spacing=2)

add_text(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(0.4),
         "Each phase includes a checkpoint review, mentor feedback, and catch-up buffer.",
         size=12, color=C_TEXT_SECONDARY, align=PP_ALIGN.CENTER, spacing=2)

slide_number(slide, 6, 19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 — 14 Topics Overview
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 6, "Project Topics",
                   "14 hands-on topics across the compiler stack, organized by difficulty")

# Group by difficulty
for diff in ["低", "中", "高"]:
    label = DIFF_LABELS[diff]
    color = DIFF_COLORS[diff]
    items = [t for t in TOPICS if t[2] == diff]

    col = {"低": 0, "中": 1, "高": 2}[diff]
    x_base = Inches(0.5 + col * 4.2)

    # Column header
    head = add_rect(slide, x_base, Inches(2.6), Inches(3.9), Inches(0.45),
                    color, radius=False)
    add_text(slide, x_base + Inches(0.2), Inches(2.63), Inches(3.5), Inches(0.4),
             f"  {label}", size=13, color=C_WHITE, bold=True, spacing=2)

    for j, (n, name, _, desc) in enumerate(items):
        y = Inches(3.2 + j * 0.28)
        add_text(slide, x_base + Inches(0.2), y, Inches(0.5), Inches(0.25),
                 n, size=9, color=color, bold=True, font="JetBrains Mono",
                 spacing=1)
        add_text(slide, x_base + Inches(0.8), y, Inches(3.0), Inches(0.25),
                 name, size=11, color=C_TEXT_PRIMARY, spacing=1)

# Separator lines between columns
for col_x in [Inches(4.5), Inches(8.7)]:
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, col_x, Inches(2.6), Inches(0.01), Inches(4.2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x25, 0x2E, 0x42)
    line.line.fill.background()

slide_number(slide, 7, 19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 8 — Beginner Topics Detail
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 7, "Beginner Track",
                   "Perfect for getting started — quick wins, real contributions")

beginner_items = [t for t in TOPICS if t[2] == "低"]
for i, (n, name, _, desc) in enumerate(beginner_items):
    x = Inches(0.5 + (i % 3) * 4.2)
    y = Inches(2.8 + (i // 3) * 2.0)

    card = add_rect(slide, x, y, Inches(3.9), Inches(1.7),
                    C_BG_CARD, radius=True, border_color=C_GREEN, border_width=Pt(0.5))

    add_tag(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.7), Inches(0.25),
            "Beginner", C_GREEN)
    add_text(slide, x + Inches(1.1), y + Inches(0.15), Inches(2.6), Inches(0.35),
             f"{n}  {name}", size=14, color=C_GREEN, bold=True, spacing=2)
    add_text(slide, x + Inches(0.3), y + Inches(0.65), Inches(3.4), Inches(0.8),
             desc, size=12, color=C_TEXT_SECONDARY, spacing=3)

slide_number(slide, 8, 19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 — Intermediate Topics Detail
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 8, "Intermediate Track",
                   "Dig deeper into the compiler core")

intermediate_items = [t for t in TOPICS if t[2] == "中"]
for i, (n, name, _, desc) in enumerate(intermediate_items):
    x = Inches(0.5 + (i % 3) * 4.2)
    y = Inches(2.8 + (i // 3) * 2.0)

    card = add_rect(slide, x, y, Inches(3.9), Inches(1.7),
                    C_BG_CARD, radius=True, border_color=C_ORANGE, border_width=Pt(0.5))

    add_tag(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.85), Inches(0.25),
            "Intermediate", C_ORANGE)
    add_text(slide, x + Inches(1.2), y + Inches(0.15), Inches(2.5), Inches(0.35),
             f"{n}  {name}", size=14, color=C_ORANGE, bold=True, spacing=2)
    add_text(slide, x + Inches(0.3), y + Inches(0.65), Inches(3.4), Inches(0.8),
             desc, size=12, color=C_TEXT_SECONDARY, spacing=3)

slide_number(slide, 9, 19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 — Advanced Topics Detail
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 9, "Advanced Track",
                   "Graduate-level compiler techniques, real algorithmic depth")

advanced_items = [t for t in TOPICS if t[2] == "高"]
for i, (n, name, _, desc) in enumerate(advanced_items):
    x = Inches(0.5 + i * 3.2)
    y = Inches(2.8)

    card = add_rect(slide, x, y, Inches(2.9), Inches(3.2),
                    C_BG_CARD, radius=True, border_color=C_RED, border_width=Pt(1))

    add_tag(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.7), Inches(0.25),
            "Advanced", C_RED)
    add_text(slide, x + Inches(0.2), y + Inches(0.6), Inches(2.5), Inches(0.5),
             f"{n}\n{name}", size=16, color=C_RED, bold=True, spacing=4)
    add_text(slide, x + Inches(0.2), y + Inches(1.4), Inches(2.5), Inches(1.5),
             desc, size=12, color=C_TEXT_SECONDARY, spacing=3)

    # Mentor note
    add_text(slide, x + Inches(0.2), y + Inches(2.8), Inches(2.5), Inches(0.3),
             "\U0001f468‍\U0001f3eb  Mentor-guided throughout",
             size=10, color=C_TEXT_SECONDARY, spacing=2)

slide_number(slide, 10, 19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 11 — Who Should Join
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 10, "Who Is This For?",
                   "No compiler background? No problem.")

quals = [
    ("\U0001f393", "Undergrads & self-taught programmers",
     "Sophomore to junior year, or equivalent self-study level"),
    ("\U0001f4bb", "Basic programming in Python or C",
     "Loops, arrays, functions — that’s enough to start"),
    ("❓", "Curious about how computers really work",
     "You wonder what happens between print() and the CPU"),
    ("\U0001f4aa", "Willing to write code and make mistakes",
     "Bugs are learning opportunities, not failures"),
    ("\U0001f552", "Can commit 8–10 hours per week",
     "Weekend binges or daily sprints — your call"),
]

for i, (emoji, title, desc) in enumerate(quals):
    y = Inches(2.5 + i * 0.85)
    card = add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.7),
                    C_BG_CARD, radius=True, border_color=C_TEAL, border_width=Pt(0.3))

    add_text(slide, Inches(1.0), y + Inches(0.1), Inches(0.5), Inches(0.5),
             emoji, size=20, color=C_WHITE, align=PP_ALIGN.CENTER, spacing=2)
    add_text(slide, Inches(1.7), y + Inches(0.08), Inches(4), Inches(0.3),
             title, size=14, color=C_WHITE, bold=True, spacing=2)
    add_text(slide, Inches(1.7), y + Inches(0.38), Inches(9), Inches(0.3),
             desc, size=12, color=C_TEXT_SECONDARY, spacing=2)

# Not required
add_text(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.3),
         "You do NOT need: compiler theory · RISC-V knowledge · assembly experience · an FPGA board",
         size=13, color=C_TEXT_SECONDARY, align=PP_ALIGN.CENTER, spacing=2)

slide_number(slide, 11, 19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 12 — Timeline
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 11, "Timeline & Milestones",
                   "From kickoff to demo day")

milestones = [
    ("Now", "Registration Open", C_TEAL),
    ("Jun 20", "Info Session &\nRegistration Close", C_TEAL_DIM),
    ("Jul 10", "Project Kickoff\nWeek 1 Launch", C_PURPLE),
    ("Aug 1", "Phase 1\nCheckpoint", C_PURPLE_DIM),
    ("Aug 28", "Phase 2\nCheckpoint", C_ORANGE),
    ("Sep 27", "Final Demo &\nWrap-up", C_RED),
]

for i, (date, event, color) in enumerate(milestones):
    x = Inches(0.5 + i * 2.15)
    y = Inches(3.0)

    # Node
    node = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(0.7), y, Inches(0.35), Inches(0.35)
    )
    node.fill.solid()
    node.fill.fore_color.rgb = color
    node.line.fill.background()

    # Connector
    if i < len(milestones) - 1:
        conn = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x + Inches(1.05), y + Inches(0.15),
            Inches(1.1), Inches(0.04)
        )
        conn.fill.solid()
        conn.fill.fore_color.rgb = color
        conn.line.fill.background()

    add_text(slide, x, y + Inches(0.5), Inches(1.9), Inches(0.3),
             date, size=13, color=color, bold=True, align=PP_ALIGN.CENTER,
             font="JetBrains Mono", spacing=2)
    add_text(slide, x, y + Inches(0.85), Inches(1.9), Inches(0.6),
             event, size=11, color=C_TEXT_SECONDARY, align=PP_ALIGN.CENTER,
             spacing=3)

# Weekly commitment
add_text(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.4),
         "\U0001f552  8–10 hours/week  ·  Weekly mentor sessions  ·  Async support on Discord",
         size=14, color=C_TEAL, align=PP_ALIGN.CENTER, spacing=2)

slide_number(slide, 12, 19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 13 — FAQ
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 12, "FAQ",
                   "Common questions answered")

faqs = [
    ("I’ve never heard of ONNX. Can I still join?",
     "Absolutely. Week 1 walks through a complete example. Think of ONNX as a “file format for models” — we’ll teach you everything."),
    ("I haven’t taken a compilers course. Will I be lost?",
     "Not at all. We use intuitive analogies (compiler as “translator”) and learn by doing, not by theory. You’ll build understanding from the ground up."),
    ("Do I need a RISC-V development board?",
     "No. Everything runs on the tinyfive software simulator on your laptop. Zero hardware required."),
    ("What if I fall behind?",
     "Every phase has a checkpoint with buffer time. Mentors check in proactively. You can focus on the core path and skip optional optimizations."),
]

for i, (q, a) in enumerate(faqs):
    y = Inches(2.5 + i * 1.1)

    q_box = add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.45),
                     C_TEAL, radius=False)
    add_text(slide, Inches(1.0), y + Inches(0.05), Inches(11), Inches(0.35),
             f"Q:  {q}", size=13, color=C_WHITE, bold=True, spacing=2)

    add_text(slide, Inches(1.0), y + Inches(0.55), Inches(11), Inches(0.45),
             f"A:  {a}", size=12, color=C_TEXT_SECONDARY, spacing=3)

slide_number(slide, 13, 19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 14 — Mentorship & Community
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 13, "Community & Mentorship",
                   "You won’t be doing this alone")

features = [
    ("\U0001f468‍\U0001f3eb", "Weekly Mentor Sessions",
     "Live deep-dives every week.\nQ&A, code reviews, and\ndiscussion of tricky concepts."),
    ("\U0001f91d", "Peer Code Review",
     "Review each other’s code.\nCatch bugs together.\nLearn by explaining your choices."),
    ("\U0001f4ac", "Always-On Community",
     "Discord / WeChat group.\nQuestions answered within hours.\nShare wins, debug together."),
    ("\U0001f3af", "Milestone Celebrations",
     "Phase checkpoints with feedback.\nFinal demo day with presentations.\nYour work becomes part of the project."),
]

for i, (emoji, title, desc) in enumerate(features):
    x = Inches(0.5 + i * 3.2)
    y = Inches(2.8)

    card = add_rect(slide, x, y, Inches(2.9), Inches(3.0),
                    C_BG_CARD, radius=True, border_color=C_TEAL, border_width=Pt(0.5))

    add_text(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.5), Inches(0.5),
             f"{emoji}", size=28, color=C_WHITE, align=PP_ALIGN.CENTER, spacing=2)
    add_text(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.5), Inches(0.4),
             title, size=16, color=C_TEAL, bold=True, align=PP_ALIGN.CENTER, spacing=2)
    add_text(slide, x + Inches(0.2), y + Inches(1.4), Inches(2.5), Inches(1.4),
             desc, size=12, color=C_TEXT_SECONDARY, align=PP_ALIGN.CENTER, spacing=3)

slide_number(slide, 14, 19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 15 — How to Join
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 14, "How to Join",
                   "Three steps to get started")

steps = [
    ("Step 1", "Register", "Fill out the registration form\nby June 20th.", C_TEAL),
    ("Step 2", "Attend Info Session", "Join us online to pick your\ntopic and meet your mentor.", C_PURPLE),
    ("Step 3", "Start Building", "Kickoff on July 10th.\nYour compiler journey begins.", C_ORANGE),
]

for i, (step, title, desc, color) in enumerate(steps):
    x = Inches(0.5 + i * 4.2)
    y = Inches(3.0)

    card = add_rect(slide, x, y, Inches(3.9), Inches(2.5),
                    C_BG_CARD, radius=True, border_color=color, border_width=Pt(1.5))

    # Step circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(1.5), y + Inches(0.2),
        Inches(0.7), Inches(0.7)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = f"0{i+1}"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = C_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "JetBrains Mono"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text(slide, x + Inches(0.3), y + Inches(1.1), Inches(3.3), Inches(0.4),
             title, size=18, color=color, bold=True, align=PP_ALIGN.CENTER, spacing=2)
    add_text(slide, x + Inches(0.3), y + Inches(1.5), Inches(3.3), Inches(0.8),
             desc, size=12, color=C_TEXT_SECONDARY, align=PP_ALIGN.CENTER, spacing=3)

add_text(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.4),
         "Registration: forms.gle/your-link  ·  Email: mentor@scratchv.dev  ·  QQ Group: xxxxxxxxx",
         size=14, color=C_TEAL, align=PP_ALIGN.CENTER, spacing=2)

slide_number(slide, 15, 19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 16 — Open Source
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 15, "Open Source Philosophy",
                   "Built in the open, for everyone")

oss_items = [
    ("\U0001f4c2", "MIT Licensed", "Free to use, modify, and share."),
    ("\U0001f500", "GitHub Workflow", "PRs, issues, discussions — standard OSS collaboration."),
    ("\U0001f4dd", "Documentation First", "Tutorials, API docs, and architecture guides."),
    ("\U0001f30d", "Community Driven", "Your contributions help future learners worldwide."),
]

for i, (emoji, title, desc) in enumerate(oss_items):
    x = Inches(0.5 + i * 3.2)
    y = Inches(2.8)

    card = add_rect(slide, x, y, Inches(2.9), Inches(2.0),
                    C_BG_CARD, radius=True, border_color=C_TEAL, border_width=Pt(0.5))
    add_text(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.5), Inches(0.4),
             f"{emoji}  {title}", size=16, color=C_WHITE, bold=True, spacing=2)
    add_text(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.5), Inches(0.8),
             desc, size=12, color=C_TEXT_SECONDARY, spacing=3)

# GitHub stats placeholder
gh_box = add_rect(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.2),
                  C_BG_CARD, radius=True, border_color=C_TEXT_SECONDARY, border_width=Pt(0.3))
add_text(slide, Inches(1.0), Inches(5.4), Inches(11), Inches(0.8),
         "\U0001f4bb  github.com/scratchv/scratchv-compiler    ·    "
         "⭐  Star us on GitHub!    ·    "
         "\U0001f504  Contributions welcome!",
         size=14, color=C_TEAL, align=PP_ALIGN.CENTER, spacing=2)

slide_number(slide, 16, 19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 17 — Team
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 16, "Mentors & Organizers",
                   "Experienced engineers guiding your journey")

mentors = [
    ("[Your Name]", "Systems / Compilers", C_TEAL),
    ("[Mentor Name]", "AI / ML Engineering", C_PURPLE),
    ("[Mentor Name]", "RISC-V / Hardware", C_ORANGE),
]

for i, (name, role, color) in enumerate(mentors):
    x = Inches(0.5 + i * 4.2)
    y = Inches(3.0)

    card = add_rect(slide, x, y, Inches(3.9), Inches(2.5),
                    C_BG_CARD, radius=True, border_color=color, border_width=Pt(0.5))

    # Avatar placeholder
    avatar = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(1.3), y + Inches(0.3),
        Inches(1.2), Inches(1.2)
    )
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = color
    avatar.fill.fore_color.brightness = 0.7
    avatar.line.fill.background()
    tf = avatar.text_frame
    tf.paragraphs[0].text = name[0]
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = C_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Inter"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text(slide, x + Inches(0.3), y + Inches(1.7), Inches(3.3), Inches(0.4),
             name, size=15, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER, spacing=2)
    add_text(slide, x + Inches(0.3), y + Inches(2.05), Inches(3.3), Inches(0.3),
             role, size=12, color=C_TEXT_SECONDARY, align=PP_ALIGN.CENTER, spacing=2)

slide_number(slide, 17, 19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 18 — CTA / Closing
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)

add_gradient_bar(slide, Inches(0), Inches(0), W, Inches(0.06), C_TEAL, C_PURPLE)

# Decorative
for x, y, s, c in [(10, -0.5, 5, C_TEAL), (8, 5.5, 4, C_PURPLE)]:
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                                    Inches(s), Inches(s))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c
    shape.fill.fore_color.brightness = 0.85
    shape.line.fill.background()

add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.0),
         "Your Compiler Journey\nStarts Today",
         size=44, color=C_WHITE, bold=True, spacing=6)

add_text(slide, Inches(0.8), Inches(3.5), Inches(10), Inches(0.5),
         "You don’t need to be great to start. You need to start to be great.",
         size=18, color=C_TEXT_SECONDARY, spacing=4)

# CTA buttons
for i, (label, sub, color) in enumerate([
    ("Register Now", "forms.gle/your-link", C_TEAL),
    ("GitHub Repo", "github.com/scratchv/scratchv-compiler", C_PURPLE),
]):
    x = Inches(0.8 + i * 3.5)
    y = Inches(4.5)

    btn = add_rect(slide, x, y, Inches(3.0), Inches(0.6),
                   color, radius=True)
    add_text(slide, x, y + Inches(0.05), Inches(3.0), Inches(0.3),
             label, size=15, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER, spacing=2)
    add_text(slide, x, y + Inches(0.3), Inches(3.0), Inches(0.25),
             sub, size=9, color=C_WHITE, align=PP_ALIGN.CENTER, spacing=1)

add_text(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.4),
         "Questions?  →  mentor@scratchv.dev  ·  QQ: xxxxxxxxx  ·  Discord: link",
         size=13, color=C_TEXT_SECONDARY, align=PP_ALIGN.CENTER, spacing=2)

add_text(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
         "#ScratchV  #Compiler  #RISC-V  #OpenSource  #LearnByBuilding",
         size=11, color=C_TEXT_SECONDARY, align=PP_ALIGN.CENTER, spacing=2)

slide_number(slide, 18, 19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 19 — Thank You
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)

add_gradient_bar(slide, Inches(0), Inches(0), W, Inches(0.06), C_TEAL, C_PURPLE)

for x, y, s, c in [(10, 0, 5, C_TEAL), (8, 4, 4, C_PURPLE), (-1, 3, 3.5, C_TEAL_DIM)]:
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                                    Inches(s), Inches(s))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c
    shape.fill.fore_color.brightness = 0.85
    shape.line.fill.background()

add_text(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(1.5),
         "Thank You\n\U0001f91d",
         size=56, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER, spacing=6)

add_text(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(0.6),
         "Let’s build something amazing together.",
         size=24, color=C_TEAL, align=PP_ALIGN.CENTER, spacing=4)

add_text(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.5),
         "Q & A",
         size=20, color=C_TEXT_SECONDARY, align=PP_ALIGN.CENTER, spacing=2)

add_text(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.4),
         "github.com/scratchv/scratchv-compiler",
         size=14, color=C_TEXT_SECONDARY, align=PP_ALIGN.CENTER, spacing=2)

add_text(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.3),
         "#ScratchV",
         size=10, color=C_TEXT_SECONDARY, align=PP_ALIGN.CENTER, spacing=2)

slide_number(slide, 19, 19)


# ── Save ──
output = "/home/kinsomwang/workspace/ScratchV/ScratchV_宣讲PPT_Pro.pptx"
prs.save(output)
print(f"✅  PPT saved: {output}")
print(f"    Slides: {len(prs.slides)}")
