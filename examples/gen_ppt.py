#!/usr/bin/env python3
"""Generate a promotional PPT for the ScratchV project."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color palette
DARK_BG = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT_BLUE = RGBColor(0x3A, 0x82, 0xF7)
ACCENT_CYAN = RGBColor(0x00, 0xd2, 0xff)
ACCENT_GREEN = RGBColor(0x00, 0xc9, 0x7a)
ACCENT_ORANGE = RGBColor(0xff, 0x6b, 0x35)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xcc, 0xcc, 0xdd)
DIM_WHITE = RGBColor(0xaa, 0xaa, 0xcc)
CARD_BG = RGBColor(0x25, 0x25, 0x45)
SECTION_BG = RGBColor(0x16, 0x16, 0x2e)


def add_bg(slide, color=DARK_BG):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, color, left, top, width, height):
    """Add a colored rectangle as background element."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """Add a text box with formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=13,
                    color=LIGHT_GRAY, font_name="Microsoft YaHei"):
    """Add a text box with multiple bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def add_card(slide, left, top, width, height, title, body, icon="",
             title_color=ACCENT_BLUE):
    """Add a card-style element with title and body."""
    # Card background
    card = add_shape_bg(slide, CARD_BG, left, top, width, height)

    # Icon + Title
    icon_text = f"{icon}  {title}" if icon else title
    add_text_box(slide, left + Inches(0.15), top + Inches(0.1),
                 width - Inches(0.3), Inches(0.4),
                 icon_text, font_size=13, color=title_color, bold=True)

    # Body
    add_text_box(slide, left + Inches(0.15), top + Inches(0.5),
                 width - Inches(0.3), height - Inches(0.6),
                 body, font_size=11, color=LIGHT_GRAY)


def add_header(slide, title, subtitle="", top=Inches(0.3)):
    """Add a consistent header with accent line."""
    # Accent line
    line = add_shape_bg(slide, ACCENT_BLUE, Inches(0.5), top,
                        Inches(0.08), Inches(0.5))
    # Title
    add_text_box(slide, Inches(0.7), top, Inches(8), Inches(0.6),
                 title, font_size=28, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.7), top + Inches(0.55), Inches(8), Inches(0.4),
                     subtitle, font_size=14, color=DIM_WHITE)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===================== SLIDE 1: Title =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK_BG)

    # Decorative top bar
    add_shape_bg(slide, ACCENT_BLUE, Inches(0), Inches(0),
                 Inches(10), Inches(0.06))

    # Title
    add_text_box(slide, Inches(1), Inches(2.0), Inches(8), Inches(1.2),
                 "ScratchV", font_size=56, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(1.5), Inches(3.0), Inches(7), Inches(0.6),
                 "From ONNX to RISC-V Assembly — A Hands-On Compiler Journey",
                 font_size=20, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide, Inches(2), Inches(3.8), Inches(6), Inches(0.8),
                 "Build your own AI model compiler from scratch in 12 weeks.\n"
                 "No prior compiler experience needed.",
                 font_size=14, color=DIM_WHITE, alignment=PP_ALIGN.CENTER)

    # Pipeline visual
    pipeline_text = "ONNX Model  →  Custom IR  →  Optimizer  →  RISC-V Assembly"
    add_text_box(slide, Inches(1), Inches(5.0), Inches(8), Inches(0.5),
                 pipeline_text, font_size=15, color=ACCENT_GREEN,
                 bold=True, alignment=PP_ALIGN.CENTER)

    # Bottom bar
    add_shape_bg(slide, ACCENT_BLUE, Inches(0), Inches(7.44),
                 Inches(10), Inches(0.06))

    # ===================== SLIDE 2: What is ScratchV =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_header(slide, "What is ScratchV?", "A minimal compiler that turns AI models into chip instructions")

    # Left column
    add_text_box(slide, Inches(0.7), Inches(1.5), Inches(4.2), Inches(0.4),
                 "🎯  The Big Idea", font_size=18, color=ACCENT_CYAN, bold=True)

    add_bullet_text(slide, Inches(0.7), Inches(2.0), Inches(4.2), Inches(3.5), [
        "Input: ONNX model (e.g., a neural network)",
        "Output: RISC-V assembly (.s file) executable on QEMU or real hardware",
        "Custom Intermediate Representation (3-address code)",
        "6 built-in optimization passes",
        "Pure Python — no LLVM/MLIR dependency",
    ])

    # Right column
    add_text_box(slide, Inches(5.5), Inches(1.5), Inches(4.2), Inches(0.4),
                 "🔬  Why It Matters", font_size=18, color=ACCENT_CYAN, bold=True)

    add_bullet_text(slide, Inches(5.5), Inches(2.0), Inches(4.2), Inches(3.5), [
        "Understand the full ML → silicon pipeline",
        "No compiler black box — every line is yours",
        "Ideal for teaching, research, and prototyping",
        "AI chip / accelerator design exploration",
        "Zero-to-one compiler construction experience",
    ])

    # Bottom highlight
    add_shape_bg(slide, CARD_BG, Inches(0.7), Inches(5.8), Inches(8.6), Inches(0.7))
    add_text_box(slide, Inches(0.9), Inches(5.85), Inches(8.2), Inches(0.6),
                 "\"You don't need to be a compiler expert to start. You just need curiosity and 8-10 hours per week.\"",
                 font_size=13, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

    # ===================== SLIDE 3: 12-Week Roadmap =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_header(slide, "12-Week Roadmap", "Structured milestones, weekly deliverables")

    phases = [
        ("W1-2", "Environment Setup", "RISC-V GCC + QEMU\nBaseline benchmarks\nONNX format basics", ACCENT_BLUE),
        ("W3-4", "IR & Parser", "Custom 3-address IR\nONNX parser (Add, Mul)\nIR text dump", ACCENT_CYAN),
        ("W5-6", "Optimization", "Constant folding\nDead code elimination\nMore ops (ReLU, GELU, MatMul)", ACCENT_GREEN),
        ("W7-8", "Backend Part I", "Instruction selection\nNaive reg allocation\nBasic block assembly", ACCENT_ORANGE),
        ("W9-10", "Backend Part II", "Greedy reg alloc\nLoop support\nBenchmark validation", RGBColor(0xa2, 0x55, 0xff)),
        ("W11-12", "Docs & Polish", "Design document\nUser manual\nFinal presentation", RGBColor(0xff, 0x41, 0xb5)),
    ]

    for i, (week, title, desc, color) in enumerate(phases):
        col = i % 3
        row = i // 3
        left = Inches(0.5 + col * 3.15)
        top = Inches(1.5 + row * 2.9)

        # Card
        card = add_shape_bg(slide, CARD_BG, left, top, Inches(2.9), Inches(2.5))
        # Top accent
        add_shape_bg(slide, color, left, top, Inches(2.9), Inches(0.06))
        # Week label
        add_text_box(slide, left + Inches(0.15), top + Inches(0.15),
                     Inches(2.6), Inches(0.3),
                     week, font_size=11, color=color, bold=True)
        # Title
        add_text_box(slide, left + Inches(0.15), top + Inches(0.4),
                     Inches(2.6), Inches(0.3),
                     title, font_size=14, color=WHITE, bold=True)
        # Description
        add_text_box(slide, left + Inches(0.15), top + Inches(0.8),
                     Inches(2.6), Inches(1.5),
                     desc, font_size=11, color=LIGHT_GRAY)

    # ===================== SLIDE 4: Architecture =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_header(slide, "Project Architecture", "Modular design, 4 core components")

    # Pipeline boxes
    boxes = [
        ("Frontend", "ONNX Parser\nDSL Parser", Inches(0.3), ACCENT_BLUE),
        ("IR", "3-Address Code\nBuilder + Printer", Inches(2.7), ACCENT_CYAN),
        ("Optimizer", "5 Passes:\nCF, DCE, Peephole\nLICM, MulAddFusion", Inches(5.1), ACCENT_GREEN),
        ("Backend", "Instr Selection\nReg Allocation\nAsm Emission", Inches(7.5), ACCENT_ORANGE),
    ]

    for i, (name, desc, left, color) in enumerate(boxes):
        # Main box
        box = add_shape_bg(slide, CARD_BG, left, Inches(1.6), Inches(2.2), Inches(1.8))
        add_shape_bg(slide, color, left + Inches(0.05), Inches(1.65), Inches(0.06), Inches(1.7))
        add_text_box(slide, left + Inches(0.2), Inches(1.7),
                     Inches(1.8), Inches(0.35),
                     name, font_size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.2), Inches(2.1),
                     Inches(1.8), Inches(1.2),
                     desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

        # Arrow between boxes
        if i < len(boxes) - 1:
            add_text_box(slide, left + Inches(2.0), Inches(2.2),
                         Inches(0.8), Inches(0.4),
                         "  ▶", font_size=20, color=DIM_WHITE, alignment=PP_ALIGN.CENTER)

    # Bottom: file tree
    add_shape_bg(slide, CARD_BG, Inches(0.5), Inches(4.0), Inches(9.0), Inches(3.0))

    tree_text = (
        "scratchv/\n"
        "├── ir/               # Core IR: types, builder, printer\n"
        "├── frontend/         # ONNX parser, DSL parser\n"
        "├── optimizer/        # 5 optimization passes\n"
        "├── backend/          # Instruction select, reg alloc, asm emit\n"
        "├── simulator/        # TinyFive adapter for verification\n"
        "├── main.py           # CLI entry point\n"
        "├── docs/             # Verification & optimization guides\n"
        "└── tests/            # 37+ unit tests"
    )
    add_text_box(slide, Inches(0.7), Inches(4.1), Inches(8.6), Inches(2.8),
                 tree_text, font_size=11, color=ACCENT_CYAN, font_name="Consolas")

    # ===================== SLIDE 5: Verification =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_header(slide, "Verification Workflow",
               "Run your generated assembly and count instructions")

    # Flow
    flow_items = [
        ("Compile", "scratchv model.onnx\n  --optimize all"),
        ("Simulate", "TinyFive / QEMU\nSpike / Renode"),
        ("Profile", "Instruction counts\nPerformance metrics"),
        ("Iterate", "Tune passes\nRecompile"),
    ]

    for i, (step, desc) in enumerate(flow_items):
        left = Inches(0.4 + i * 2.45)
        box = add_shape_bg(slide, CARD_BG, left, Inches(1.6), Inches(2.2), Inches(1.8))
        add_shape_bg(slide, ACCENT_BLUE, left + Inches(0.05), Inches(1.65),
                     Inches(0.06), Inches(1.7))

        add_text_box(slide, left + Inches(0.2), Inches(1.7),
                     Inches(1.8), Inches(0.3),
                     f"0{i+1}", font_size=24, color=ACCENT_BLUE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.2), Inches(2.0),
                     Inches(1.8), Inches(0.3),
                     step, font_size=15, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.2), Inches(2.4),
                     Inches(1.8), Inches(0.9),
                     desc, font_size=11, color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)

        if i < len(flow_items) - 1:
            add_text_box(slide, left + Inches(2.15), Inches(2.2),
                         Inches(0.4), Inches(0.4),
                         "→", font_size=24, color=DIM_WHITE,
                         alignment=PP_ALIGN.CENTER)

    # Tools table
    tools = (
        "TinyFive    Pure Python RV32IM simulator      pip install tinyfive\n"
        "Spike       RISC-V official ISA simulator      riscv-isa-sim\n"
        "QEMU        Industrial system emulator         apt install qemu-user\n"
        "Renode      Embedded system simulator          renode.io"
    )
    add_shape_bg(slide, CARD_BG, Inches(0.5), Inches(4.0), Inches(9.0), Inches(1.5))
    add_text_box(slide, Inches(0.5), Inches(3.8), Inches(9.0), Inches(0.3),
                 "🔧  Supported Simulators", font_size=14, color=ACCENT_CYAN, bold=True)
    add_text_box(slide, Inches(0.7), Inches(4.2), Inches(8.6), Inches(1.2),
                 tools, font_size=12, color=LIGHT_GRAY, font_name="Consolas")

    # Bottom CTA
    add_shape_bg(slide, CARD_BG, Inches(0.5), Inches(5.8), Inches(9.0), Inches(0.7))
    add_text_box(slide, Inches(0.7), Inches(5.85), Inches(8.6), Inches(0.6),
                 "💡  Measure optimization impact: compare instruction counts before vs. after",
                 font_size=13, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

    # ===================== SLIDE 6: Optimization Passes =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_header(slide, "Optimization Passes",
               "6 beginner-friendly passes — implement one per week")

    passes = [
        ("常量折叠\nConstant Folding", "Compile-time constant\nevaluation", "⭐"),
        ("死代码消除\nDead Code Elim.", "Remove unused\ninstructions", "⭐⭐"),
        ("Mul-Add Fusion", "Combine mul+add\nto reduce regs", "⭐"),
        ("窥孔优化\nPeephole", "Eliminate redundant\npatterns", "⭐"),
        ("循环不变代码外提\nLICM", "Hoist invariants\nout of loops", "⭐⭐"),
        ("贪心寄存器分配\nGreedy Reg Alloc", "LRU-based alloc\nreduce spilling", "⭐⭐"),
    ]

    for i, (name, desc, diff) in enumerate(passes):
        col = i % 3
        row = i // 3
        left = Inches(0.5 + col * 3.15)
        top = Inches(1.5 + row * 2.7)

        card = add_shape_bg(slide, CARD_BG, left, top, Inches(2.9), Inches(2.3))
        add_shape_bg(slide, ACCENT_GREEN, left, top, Inches(0.06), Inches(2.3))

        add_text_box(slide, left + Inches(0.2), top + Inches(0.15),
                     Inches(2.5), Inches(0.7),
                     name, font_size=12, color=WHITE, bold=True)
        add_text_box(slide, left + Inches(0.2), top + Inches(0.85),
                     Inches(2.5), Inches(0.8),
                     desc, font_size=11, color=LIGHT_GRAY)
        add_text_box(slide, left + Inches(0.2), top + Inches(1.7),
                     Inches(2.5), Inches(0.3),
                     f"Difficulty: {diff}", font_size=10, color=DIM_WHITE)

    # ===================== SLIDE 7: Target Audience =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_header(slide, "Who Is This For?", "No compiler expertise required")

    audiences = [
        ("🎓", "Students", "CS / EE undergrads\nWant to understand\ncompilers & AI", ACCENT_BLUE),
        ("🔬", "Researchers", "AI chips / accelerators\nNeed rapid prototyping\nCustom ISA exploration", ACCENT_CYAN),
        ("💻", "Self-taught Devs", "Curious about \"how code\nruns on silicon\"\nHands-on learners", ACCENT_GREEN),
        ("🏫", "Educators", "Compiler design course\nProject-based teaching\nOpen-source materials", ACCENT_ORANGE),
    ]

    for i, (icon, title, desc, color) in enumerate(audiences):
        left = Inches(0.5 + i * 2.4)
        card = add_shape_bg(slide, CARD_BG, left, Inches(1.6), Inches(2.15), Inches(2.8))
        add_shape_bg(slide, color, left, Inches(1.6), Inches(2.15), Inches(0.06))

        add_text_box(slide, left, Inches(1.8), Inches(2.15), Inches(0.5),
                     icon, font_size=32, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(2.3), Inches(2.15), Inches(0.3),
                     title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.15), Inches(2.7),
                     Inches(1.85), Inches(1.5),
                     desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Prerequisites
    add_shape_bg(slide, CARD_BG, Inches(0.5), Inches(4.8), Inches(9.0), Inches(2.0))
    add_text_box(slide, Inches(0.7), Inches(4.9), Inches(8.6), Inches(0.3),
                 "📋  Prerequisites", font_size=14, color=ACCENT_CYAN, bold=True)
    add_bullet_text(slide, Inches(0.7), Inches(5.3), Inches(8.6), Inches(1.3), [
        "Basic Python or C programming (variables, loops, functions)",
        "8-10 hours per week commitment",
        "No compiler theory required — we teach it from the ground up",
        "No RISC-V knowledge needed — you'll learn it in weeks 1-2",
    ])

    # ===================== SLIDE 8: Example Code =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_header(slide, "See It In Action", "From 3 lines of DSL to RISC-V assembly")

    # Code side by side
    # Left: DSL
    add_shape_bg(slide, CARD_BG, Inches(0.5), Inches(1.5), Inches(4.3), Inches(3.0))
    add_text_box(slide, Inches(0.7), Inches(1.55), Inches(3.9), Inches(0.3),
                 "📝  DSL Input", font_size=14, color=ACCENT_CYAN, bold=True)
    dsl_code = (
        "# ReLU activation\n"
        "t1 = add(input, bias)\n"
        "y = relu(t1)\n"
        "return y"
    )
    add_text_box(slide, Inches(0.7), Inches(1.9), Inches(3.9), Inches(2.4),
                 dsl_code, font_size=13, color=ACCENT_GREEN, font_name="Consolas")

    # Right: Assembly output
    add_shape_bg(slide, CARD_BG, Inches(5.2), Inches(1.5), Inches(4.3), Inches(3.0))
    add_text_box(slide, Inches(5.4), Inches(1.55), Inches(3.9), Inches(0.3),
                 "⚙️  RISC-V Output", font_size=14, color=ACCENT_ORANGE, bold=True)
    asm_code = (
        ".globl main\n"
        "main:\n"
        "  add t2, t0, t1\n"
        "  max t3, t2, x0\n"
        "  mv   a0, t3\n"
        "  ret"
    )
    add_text_box(slide, Inches(5.4), Inches(1.9), Inches(3.9), Inches(2.4),
                 asm_code, font_size=13, color=ACCENT_ORANGE, font_name="Consolas")

    # Bottom: Pipeline
    add_shape_bg(slide, CARD_BG, Inches(0.5), Inches(4.8), Inches(9.0), Inches(1.2))
    add_text_box(slide, Inches(0.7), Inches(4.9), Inches(8.6), Inches(0.3),
                 "🔁  Pipeline: DSL → IR → Optimize → Assembly → Verify",
                 font_size=13, color=WHITE, bold=True)
    pipeline_steps = (
        "$ scratchv examples/relu_test.dsl -o relu.s --optimize all\n"
        "$ python examples/verify_with_tinyfive.py examples/relu_test.dsl\n"
        "  Instructions before: 3    Instructions after: 3    Reduction: 0.0%"
    )
    add_text_box(slide, Inches(0.7), Inches(5.25), Inches(8.6), Inches(0.6),
                 pipeline_steps, font_size=11, color=ACCENT_CYAN, font_name="Consolas")

    # ===================== SLIDE 9: Get Involved =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)

    # Decorative top
    add_shape_bg(slide, ACCENT_BLUE, Inches(0), Inches(0), Inches(10), Inches(0.06))

    # Main CTA
    add_text_box(slide, Inches(1), Inches(1.5), Inches(8), Inches(0.8),
                 "Get Involved", font_size=42, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(2.3), Inches(6), Inches(0.6),
                 "Start building your compiler today",
                 font_size=18, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

    # Info boxes
    boxes_data = [
        ("📖", "Read the Docs", "docs/verification.md\ndocs/optimization_guide.md"),
        ("💻", "Explore the Code", "github.com/scratchv\n(open source, MIT license)"),
        ("🚀", "Quick Start", "git clone && cd ScratchV\npython3 -m venv .venv && source .venv/bin/activate\npip install -e ."),
        ("🧪", "Run the Tests", "pytest tests/ -v  # 37+ tests"),
    ]

    for i, (icon, title, desc) in enumerate(boxes_data):
        col = i % 2
        row = i // 2
        left = Inches(0.8 + col * 4.7)
        top = Inches(3.2 + row * 1.7)

        card = add_shape_bg(slide, CARD_BG, left, top, Inches(4.2), Inches(1.4))
        add_shape_bg(slide, ACCENT_CYAN, left, top, Inches(4.2), Inches(0.04))

        add_text_box(slide, left + Inches(0.2), top + Inches(0.15),
                     Inches(0.5), Inches(0.4),
                     icon, font_size=24)
        add_text_box(slide, left + Inches(0.7), top + Inches(0.15),
                     Inches(3.3), Inches(0.3),
                     title, font_size=15, color=WHITE, bold=True)
        add_text_box(slide, left + Inches(0.7), top + Inches(0.5),
                     Inches(3.3), Inches(0.8),
                     desc, font_size=11, color=LIGHT_GRAY, font_name="Consolas")

    # Bottom tagline
    add_text_box(slide, Inches(1.5), Inches(6.5), Inches(7), Inches(0.5),
                 "You don't need to be great to start, but you need to start to be great.",
                 font_size=14, color=DIM_WHITE, alignment=PP_ALIGN.CENTER)

    add_shape_bg(slide, ACCENT_BLUE, Inches(0), Inches(7.44), Inches(10), Inches(0.06))

    return prs


def main():
    output_dir = "/home/kinsomwang/workspace/ScratchV"
    output_path = os.path.join(output_dir, "ScratchV_Promo.pptx")

    prs = create_presentation()
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    print(f"   Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
